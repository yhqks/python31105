# -*- coding: utf-8 -*-
import os
from pathlib import Path
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter
import io
from langsmith import expect
import json
from typing import List, Tuple,Union
import tempfile
from paddleocr import PaddleOCR,PPStructure
from collections import defaultdict
from typing import Union, List, Dict
import logging


# ========== 模块导入 ==========
try:
    from docx import Document
    from pptx import Presentation
    from docx.parts.image import ImagePart
    from xlrd import xldate_as_datetime
    import xlrd
    from openpyxl import load_workbook
    import fitz  # PyMuPDF
    import pandas as pd
    import pdfplumber
    # from langchain_community.vectorstores import Chroma
    # from openai import OpenAI
    import chardet
except ImportError:
    raise ImportError("请先安装依赖库：python-docx, python-pptx, openpyxl, pdfplumber,Chroma,openai,chardet，docx,PaddleOCR,PPStructure,xlrd,fitz")
logger=logging.getLogger('file_parser_logger4')
logger.setLevel(logging.ERROR)
# 创建一个handler，用于写入日志文件
file_handler = logging.FileHandler('file_parser_logger4.log',encoding='utf-8')
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
file_handler.setFormatter(formatter)
logger.addHandler(file_handler)
 
# 添加handler到logger
logger.addHandler(file_handler)
class DocumentParseError(Exception):
    """自定义文档解析异常"""
    def __init__(self, message: str,file_path: str = None) -> None:
        super().__init__(message)
        self.res(message,file_path)
    def res(self,message,file_path)->None:
        # 这里可以添加日志记录或其他处理逻辑
        logger.error(f"解析失败：{message}\nfile_path:{file_path}")
       



# ========== 核心解析器 ==========
class UniversalDocumentParser:
    """通用文档解析器"""
    def __init__(self):
        self.supported_formats = {
            '.docx': self._parse_word,
            '.pptx': self._parse_ppt,
            '.xlsx': self._parse_excel,
            '.csv': self._parse_excel,
            '.pdf': self._parse_pdf,
            # '.doc': self._parse_word,  # 支持Word 97-2003文档
            '.txt': self._parse_txt,
            '.xls': self._parse_excel,  # 支持Excel 97-2003文档
            '.jpg': self._ocr_image,
            '.png': self._ocr_image,
            '.jpeg': self._ocr_image,
            '.bmp': self._ocr_image,
            '.gif': self._ocr_image,
            '.tiff': self._ocr_image,
        }
        self.table_engine = PPStructure(show_log=False)  # 初始化表格识别模
        # 初始化常规OCR模型（用于文本识别）
        self.text_engine = PaddleOCR(use_angle_cls=True, lang="ch")
        self.temp_dir = tempfile.TemporaryDirectory(dir='D:\yhq\python31105/temp')  # 用于暂存提取的图片
    def __del__(self):
        self.temp_dir.cleanup()  # 清理临时文件

    def parse(self, file_path: str) -> str:

        """
        通用文档解析入口
        :param file_path: 文件路径
        :return: 提取的文本内容
        """
        # 校验文件是否存在
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在：{file_path}")

        # 获取文件扩展名
        ext = os.path.splitext(file_path)[1].lower()
        print(f"文件扩展名：{ext}")

        # 选择解析方法
        parser = self.supported_formats.get(ext)
        print(f"选择的解析方法：{parser}")
        if not parser:
                raise DocumentParseError(f"不支持的文件格式：{ext}",file_path=file_path)

        try:
            return parser(file_path)
        except Exception as e:
            raise DocumentParseError(f"解析失败：{str(e)}",file_path=file_path) from e

    def _initialize_ocr(self):
        if self.ocr_model is None:
            self.ocr_model = PaddleOCR(
                use_angle_cls=True,
                lang='ch',
                table=True,
                  show_log=False,
            )



    def _ocr_image(self, image_path: str) -> Union[str, Dict]:
        """结构化OCR识别（表格优先）"""
        try:
    
            # 使用PPStructure处理图片
            table_result = self.table_engine(image_path)
          
            # 提取所有表格区域
            tables = [item for item in table_result if item['type'] == 'table']
            
            if tables:
                
                # 取第一个有效表格
                first_table = tables[0]
                return self._parse_table_data(first_table)
                
            else:
                # 没有表格时使用常规OCR识别文本
        
                text_result = self.text_engine.ocr(image_path, cls=True)
                return self._parse_text_data(text_result)
                
        except Exception as e:
             raise DocumentParseError(f"图片ocr失败：{str(e)}",file_path=image_path) from e
        
    def _parse_table_data(self, table_item: Dict) -> Dict:
        """解析表格数据结构为键值对格式（适配PPStructure输出）"""
        try:
            # 验证输入结构
            if table_item.get('type') != 'table':
                return {"error": "非表格类型数据", "actual_type": table_item.get('type')}
            
            res_data = table_item.get('res', {})
            if not res_data:
                return {"error": "缺失表格解析结果", "input": table_item}
            
            # 从HTML中提取结构化数据（更可靠）
            html_content = res_data.get('html', '')
            if not html_content:
                return {"error": "缺失表格HTML内容"}
            
            # 解析HTML表格
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            table = soup.find('table')
            if not table:
                return {"error": "HTML中未找到表格"}

            # 提取表头和行数据
            rows = []
            for tr in table.find_all('tr'):
                cells = [td.get_text(strip=True) for td in tr.find_all('td')]
                if cells:  # 忽略空行
                    rows.append(cells)
            
            if len(rows) < 2:
                return {"type": "table", "warning": "无效表格结构（行数不足）"}
            
            # 处理表头（合并多行表头）
            headers = []
            for i, header in enumerate(rows[0]):
                if not header:
                    header = f"Column_{i+1}"
                headers.append(header)
            
            # 构建键值对
            table_data = []
            for row in rows[1:]:
                row_dict = {}
                for i in range(len(headers)):
                    # 处理列数不匹配的情况
                    value = row[i] if i < len(row) else ""
                    row_dict[headers[i]] = value
                table_data.append(row_dict)
            
            return {
                "type": "table",
                "format": "header-row mapping",
                "headers": headers,
                "rows": table_data,
                "dimension": f"{len(table_data)}行x{len(headers)}列",
                "confidence": table_item.get('score', 0),
                "source": "PPStructure HTML"
            }
            
        except Exception as e:
            return {
                "error": "表格解析失败",
                "exception": str(e),
                "input_snapshot": {
                    "type": table_item.get('type'),
                    "score": table_item.get('score'),
                    "html_preview": str(table_item.get('res', {}).get('html', '')[:200] + "...")
                }
            }
        

    def _parse_text_data(self, ocr_result: list) -> str:
        """解析常规文本结果"""
        text_lines = []
        for region in ocr_result:
            if region is None:
                continue
            for line in region:
                if isinstance(line, (list, tuple)) and len(line) >= 2:
                    text = str(line[1][0]).strip()
                    if text:
                        text_lines.append(text)
        return '\n'.join(text_lines) if text_lines else "未识别到文本"
    

    def _format_table_output(self, table_cells: list) -> str:
        table_data = {
            "type": "table",
            "data": [[cell['text'] for cell in row] for row in table_cells]
        }
        return json.dumps(table_data, ensure_ascii=False, indent=2)
        
    def _get_image_extension(self, part: ImagePart) -> str:
        """获取图片的扩展名"""
        # 根据图片的MIME类型获取扩展名
        mime_type = part.content_type
        if mime_type == "image/png":
            return "png"
        elif mime_type == "image/jpeg":
            return "jpg"
        elif mime_type == "image/gif":
            return "gif"
        elif mime_type == "image/bmp":
            return "bmp"
        elif mime_type == "image/tiff":
            return "tiff"
        else:
            return "bin"  # 默认扩展名


    def _extract_images_from_ppt(self, prs: Presentation) -> List[Tuple[str, str]]:
        """从PPT提取图片"""
        images = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # 图片类型
                    image_data = shape.image.blob
                    ext = shape.image.ext
                    with tempfile.NamedTemporaryFile(
                        dir=self.temp_dir.name,
                        suffix=f".{ext}",
                        delete=False
                    ) as f:
                        preprocess_image=self.preprocess_image(image_data)
                        f.write(preprocess_image)

                        images.append(f.name)
        return [(img, self._ocr_image(img)) for img in images]
    
    def preprocess_image(binary_data):
        # 将二进制数据转换为图像
        image = Image.open(io.BytesIO(binary_data))
        
        # 1. 提升图像质量
        # 调整分辨率（如果需要）
        image = image.resize((image.width * 2, image.height * 2), Image.LANCZOS)
        
        # 增强对比度
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(1.5)  # 增强对比度
        
        # 2. 去噪与校正
        # 去噪（中值滤波）
        image = image.filter(ImageFilter.MedianFilter(size=3))
        
        # 文字方向校正（使用pytesseract检测文字方向）
        try:
            osd = pytesseract.image_to_osd(image)
            angle = int(osd.split('\n')[1].split(':')[1].strip())
            if angle != 0:
                image = image.rotate(angle, expand=True)
        except Exception as e:
            print(f"文字方向校正失败: {e}")
        
        # 3. 灰度转换与二值化
        # 转换为灰度图像
        image = image.convert('L')
        
        # 二值化处理
        threshold = 128
        image = image.point(lambda p: 255 if p > threshold else 0)
        
        # 将处理后的图像转换为二进制数据
        output_buffer = io.BytesIO()
        image.save(output_buffer, format='PNG')
        processed_binary_data = output_buffer.getvalue()
        
        return processed_binary_data

    def _extract_images_from_pdf(self, file_path: str) -> dict[int, list[tuple[str, str]]]:
        """按页码提取图片及其OCR结果"""
        page_images = defaultdict(list)
        doc = fitz.open(file_path)
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                if base_image["ext"] in ["png", "jpeg", "jpg"] and base_image.get("image"):
                    try:
                        with tempfile.NamedTemporaryFile(
                            dir=self.temp_dir.name,
                            suffix=f".{base_image['ext']}",
                            delete=False
                        ) as f:
                            f.write(base_image['image'])
                            # 不知道为啥反正f.write不能直接保存需要用image在保存一次 看不懂
                            image_bytes = base_image["image"]  # 图像的二进制数据
                            # 使用Pillow库保存图像
                            image = Image.open(io.BytesIO(image_bytes))  # 将二进制数据加载为图像
                            image.save(f.name)  # 保存图像，指定文件名和格式
                            img_path = f.name

                            # 直接执行OCR并存储结果
                            ocr_text = self._ocr_image(img_path)
                            # 按页码分组存储 (PyMuPDF页码从0开始)
                            page_images[page_num].append((img_path, ocr_text))
                    except Exception as e:
                        raise DocumentParseError(f"图片处理失败页码{page_num}:{str(e)}",file_path=file_path) from e
        return dict(page_images)
    

    def _parse_word(self, file_path: str) -> str:
        """解析Word文档，按原始顺序提取文字、表格和图片中的文字"""
        doc = Document(file_path)
        text = []
        # 创建临时目录
        with tempfile.TemporaryDirectory() as temp_dir:
            # 预构建元素映射
            element_to_para = {p._element: p for p in doc.paragraphs}
            element_to_table = {t._element: t for t in doc.tables}
            
            # 遍历文档所有元素
            for element in doc.element.body.iterchildren():
                # 处理段落
                if element in element_to_para:
                    para = element_to_para[element]
                    if para.text.strip():
                        text.append(para.text)
                    
                    # 处理段落中的图片
                    for run in para.runs:
                        for drawing in run._element.xpath(".//*[local-name()='drawing']"):
                            for pic in drawing.xpath(".//*[local-name()='pic']"):
                                blip = pic.xpath(".//*[local-name()='blip']")
                                if blip:
                                    embed_id = blip[0].get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                                    if embed_id and embed_id in doc.part.related_parts:
                                        # 提取图片内容
                                        image_part = doc.part.related_parts[embed_id]
                                        img_path = os.path.join(temp_dir, f"{embed_id}.png")
                                        with open(img_path, 'wb') as f:
                                            f.write(image_part.blob)
                                        
                                        # 调用OCR方法
                                        ocr_text = self._ocr_image(img_path)
                                        if ocr_text:
                                            text.append(f"[图片内容] {ocr_text}")
                                        
                                        # 删除临时文件
                                        os.remove(img_path)

                # 处理表格
                elif element in element_to_table:
                    table = element_to_table[element]
                    if len(table.rows) > 0:  # 确保表格有行
                        # 获取表头（第一行）
                        headers = [cell.text.strip() for cell in table.rows[0].cells]
                        
                        table_data = []
                        # 处理数据行（从第二行开始）
                        for row in table.rows[1:]:
                            row_data = {}
                            for i, cell in enumerate(row.cells):
                                if i < len(headers):  # 确保不超过表头数量
                                    row_data[headers[i]] = cell.text.strip()
                            if row_data:  # 非空行才添加
                                table_data.append(row_data)
                        text.append(json.dumps(
                        {"table": table_data},
                        ensure_ascii=False,
                        indent=2
                    ))
        return '\n'.join(text)

    def _parse_ppt(self, file_path: str) -> str:
        """解析PowerPoint文档"""
        prs = Presentation(file_path)
        text = []
        
        # 遍历所有幻灯片
        for slide in prs.slides:
            # 提取形状中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
                # 处理表格（PPTX中的表格）
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text for cell in row.cells]
                        text.append('\t'.join(row_text))
                if shape.shape_type == 13:  # 图片类型
                    image_data = shape.image.blob
                    ext = shape.image.ext
                    try:
                        with tempfile.NamedTemporaryFile(
                            dir=self.temp_dir.name,
                            suffix=f".{ext}",
                            delete=False
                        ) as f:
                            f.write(image_data)
                            image_bytes = image_data  # 图像的二进制数据
                                # 使用Pillow库保存图像
                            image = Image.open(io.BytesIO(image_bytes))  # 将二进制数据加载为图像
                            image.save(f.name)  # 保存图像，指定文件名和格式
                            text.append(self._ocr_image(f.name))
                    except Exception as e:
                        raise DocumentParseError(f"图片处理失败:{str(e)}",file_path=file_path) from e   
        # image_results = self._extract_images_from_ppt(prs)
        # for img_path, ocr_text in image_results:
        #     text.append(ocr_text)
        self.__del__()
        return ''.join(json.dumps(text,ensure_ascii=False, indent=2))
        # return  ''.join(map(str, text))
        

    def _parse_excel(self, file_path: Union[str, Path]) -> str:
        """解析Excel为结构化JSON，支持大文件与合并单元格"""
        result = {}
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise DocumentParseError("文件不存在",file_path=file_path)

        # 提取文件名信息
        file_name = file_path.stem
        result["file_name"] = file_name

        try:
            if file_path.suffix == '.xls':
                self._parse_xls(file_path, result)
            elif file_path.suffix == '.xlsx':
                self._parse_xlsx(file_path, result)
            elif file_path.suffix == '.csv':
                self._parse_csv(file_path, result)
            else:
                raise DocumentParseError("不支持的格式",file_path=file_path)
        except Exception as e:
            raise DocumentParseError(str(e),file_path=file_path) from e

        return json.dumps(result, ensure_ascii=False, indent=2)

    def _parse_xls(self, file_path, result):
        """处理.xls文件"""
        wb = xlrd.open_workbook(file_path)
        for sheet_name in wb.sheet_names():
            sheet = wb.sheet_by_name(sheet_name)
            header_row = self._detect_header(sheet, 'xls')
            headers = self._get_headers(sheet, header_row, 'xls')
            data = []
            
            for row_idx in range(header_row+1, sheet.nrows):
                row_data = {}
                for col_idx, header in enumerate(headers):
                    cell = sheet.cell(row_idx, col_idx)
                    value = self._parse_xls_cell(cell, wb)
                    row_data[header] = value
                data.append(row_data)
            
            result[sheet_name] = data

    def _parse_xlsx(self, file_path, result):
        """处理.xlsx文件（关闭只读模式以支持合并单元格）"""
        wb = load_workbook(file_path, read_only=False, data_only=True)  # 关键修改
        for sheet in wb:
            headers = self._get_headers(sheet, 0, 'xlsx')  # 假设表头在第一行
            data = []
            
            # 处理合并单元格
            merged_values = {}
            for merge_range in sheet.merged_cells.ranges:
                min_row, min_col, max_row, max_col = merge_range.min_row, merge_range.min_col,merge_range.max_row,merge_range.max_col
                master_value = sheet.cell(min_row, min_col).value
                for row in range(min_row, max_row+1):
                    for col in range(min_col, max_col+1):
                        merged_values[(row, col)] = master_value
            
            # 提取数据
            for row_idx, row in enumerate(sheet.iter_rows(min_row=2, values_only=False), start=2):
                row_data = {}
                for col_idx, cell in enumerate(row, start=1):
                    value = merged_values.get((row_idx, col_idx), cell.value)
                    header = headers[col_idx-1] if (col_idx-1) < len(headers) else f"Col{col_idx}"
                    row_data[header] = value if value is not None else ""
                data.append(row_data)
            
            result[sheet.title] = data
        wb.close()  # 显式关闭工作簿

    def _parse_csv(self, file_path, result):
        """流式处理大CSV文件"""
        chunk_size = 10**4  # 1万行/块
        data = []
        
        try:
            for chunk in pd.read_csv(file_path, chunksize=chunk_size,
                                     dtype=str, keep_default_na=False):
                data.extend(chunk.to_dict('records'))
        except pd.errors.ParserError as e:
            raise DocumentParseError(f"CSV解析错误: {str(e)}",file_path=file_path) from e
        
        result[file_path.stem] = data

    def _parse_xls_cell(self, cell, wb):
        """处理.xls单元格类型"""
        if cell.ctype == xlrd.XL_CELL_DATE:
            return xldate_as_datetime(cell.value, wb.datemode).isoformat()
        return cell.value if cell.value != "" else ""

    def _handle_merged_cells(self, sheet):
        """处理合并单元格（适用于openpyxl）"""
        for merge_range in sheet.merged_cells.ranges:
            min_row, min_col, max_row, max_col = merge_range.min_row,  merge_range.min_col,merge_range.max_row,merge_range.max_col
            master_value = sheet.cell(min_row, min_col).value
            for row in sheet.iter_rows(min_row, max_row, min_col, max_col):
                for cell in row:
                    if cell.coordinate != sheet.cell(min_row, min_col).coordinate:
                        cell.value = master_value

    def _detect_header(self, sheet, file_type) -> int:
        """智能检测表头行（示例实现）"""
        # 可扩展更复杂的检测逻辑
        return 0  # 默认识别首行为表头

    def _get_headers(self, sheet, header_row, file_type):
        """提取有效表头"""
        if file_type == 'xls':
            return [
                str(cell.value).strip() 
                for cell in sheet.row(header_row) 
                if cell.value
            ]
        else:
            return [
                str(cell.value).strip() 
                for cell in sheet[header_row+1] 
                if cell.value
            ]

    def _parse_pdf(self, file_path: str) -> str:
        """按页合并文本、表格和图片内容"""
        full_text = []
        # 先提取所有页面的图片信息 (页码 -> 图片列表)
        page_images = self._extract_images_from_pdf(file_path)

        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                page_content = []
                
                # 1. 提取文本
                if (page_text := page.extract_text()):
                    page_content.append(page_text)
                
                # 2. 提取表格
                for table in page.extract_tables():
                    if table:  # 确保表格非空
                        # 获取表头（假设第一行是表头）
                        headers = [
                            str(cell).replace('\n', ' ').strip() if cell else "" 
                            for cell in table[0]
                        ]
                        
                        table_data = []
                        # 处理数据行（从第二行开始）
                        for row in table[1:]:
                            row_data = {}
                            for i, cell in enumerate(row):
                                # 确保不超过表头数量
                                if i < len(headers):
                                    cleaned_cell = str(cell).replace('\n', ' ').strip() if cell else ""
                                    row_data[headers[i]] = cleaned_cell
                            # 只添加非空行
                            if any(row_data.values()):
                                table_data.append(row_data)
                        
                        # 将表格数据转为JSON字符串
                        if headers and table_data:  # 确保有表头和数据
                            try:
                                table_str = json.dumps(
                                    {"table_data": table_data},
                                    ensure_ascii=False,
                                    indent=2
                                )
                                page_content.append(table_str)
                            except Exception as e:
                                # 如果JSON转换失败，使用简单字符串表示
                                page_content.append(str(table_data))
                
                # 3. 添加本页图片OCR结果
                if img_pairs := page_images.get(page_num):
                    for img_path, ocr_text in img_pairs:
                        page_content.append(ocr_text)
                
                # 合并本页所有内容，确保所有元素都是字符串
                if page_content:
                    full_text.append("\n".join(str(item) for item in page_content))
    
        self.__del__()
        return "\n".join(full_text)

    def _parse_txt(self, file_path: str) -> str:
        """解析文本文件"""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

class FileEncoder:
    def __init__(self, input_file):
        """
        初始化 FileEncoder 类。
        
        :param input_file: 输入文件路径
        """
        self.input_file = input_file
        self.current_encoding = None

    def detect_encoding(self):
        """
        检测文件的当前编码格式。
        
        :return: 文件的编码格式
        """
        with open(self.input_file, "rb") as file:
            raw_data = file.read()
            result = chardet.detect(raw_data)
            self.current_encoding = result["encoding"]
            return self.current_encoding

    def convert_encoding(self, output_file, target_encoding):
        """
        将文件从当前编码转换为目标编码。
        
        :param output_file: 输出文件路径
        :param target_encoding: 目标编码格式
        """
        if not self.current_encoding:
            self.detect_encoding()

        with open(self.input_file, "r", encoding=self.current_encoding, errors="replace") as file:
            content = file.read()

        with open(output_file, "w", encoding=target_encoding, errors="replace") as file:
            file.write(content)

        print(f"文件已从 {self.current_encoding} 转换为 {target_encoding}，并保存到 {output_file}")

    def convert_to_utf8(self, output_file):
        self.convert_encoding(output_file, "utf-8")

    def convert_to_gbk(self, output_file):
        self.convert_encoding(output_file, "gbk")

    def convert_to_latin1(self, output_file):
        self.convert_encoding(output_file, "latin1")

    def convert_to_windows1252(self, output_file):
        self.convert_encoding(output_file, "windows-1252")

    def convert_to_iso_8859_1(self, output_file):
        self.convert_encoding(output_file, "iso-8859-1")


if __name__ == "__main__":
    parser = UniversalDocumentParser()

    # parser = UniversalDocumentParser()
    # # result_json = parser._ocr_image("E:/5f7efdfb-60d1-4fc4-af7e-8bc618025290.png")
    # # result_data = json.loads(result_json)  # 解析为Python对象
    test_result = parser.parse('D:/AIOPS资料及代码/AIOPS资料及代码/资料/设计院-无线智能维护AIOps【主要功能交流材料】（202402）.pdf')
    print(test_result)